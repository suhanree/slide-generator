import time
import os
import base64
import json
import PIL
import io

from pptx import Presentation
from pptx.util import Inches, Pt
from google.cloud import storage
import pandas as pd
from anthropic import AnthropicVertex
from selenium import webdriver
from selenium.webdriver.common.by import By

from langchain_google_vertexai.model_garden import ChatAnthropicVertex
from langchain_core.messages import HumanMessage, SystemMessage

MODEL_NAME = "claude-3-7-sonnet@20250219"
MAX_TOKENS = 2048
TEMPERATURE = 0.0
TOP_P = 0.95
TOP_K = 40

MAX_RETRIES = 5



INPUT_COLUMNS = ["segment_id", "feature_description", "segment_mean", "overall_mean"]

DEFAULT_IMAGE_FILENAME = "default_image"
DEFAULT_IMAGE_FILETYPE = "jpeg"
with open(f"{DEFAULT_IMAGE_FILENAME}.{DEFAULT_IMAGE_FILETYPE}", "rb") as image_file:
    DEFAULT_IMAGE_BASE64 = base64.b64encode(image_file.read()).decode("utf-8")


def query_llm(prompt, model_name=MODEL_NAME, client=None, max_tokens=MAX_TOKENS, temperature=TEMPERATURE, top_p=TOP_P, top_k=TOP_K,
              additional_contents=None):
    assert client is not None
    content = [
        {
            "type": "text",
            "text": prompt,
        },
    ]
    if additional_contents is not None:
        content.extend(additional_contents)
    response = client.messages.create(
        max_tokens=max_tokens,
        temperature=temperature,
        top_p=top_p,
        top_k=top_k,
        messages=[
            {
                "role": "user",
                "content": content,
            }
        ],
        model=model_name
    )
    return response.content[0].text


def json_parser(text):
    first = text.rfind("{")
    last = text.rfind("}")
    return json.loads(text[first:last+1])


def get_name_and_description(
    segment_id: int,
    segment_input: str,
    prompt_template1: str,
    client=None,
    model_name=MODEL_NAME, max_tokens=MAX_TOKENS, temperature=TEMPERATURE, top_p=TOP_P, top_k=TOP_K,
) -> dict:
    assert client is not None
    prompt1 = prompt_template1.format(segment_input=segment_input)
    try:
        response = query_llm(prompt1, model_name=model_name, client=client, max_tokens=max_tokens, temperature=temperature, top_p=top_p, top_k=top_k)
        return json_parser(response)
    except Exception as e:
        print(e)
        return {"segment_id": segment_id, "segment_name": "NA", "segment_description": "NA", "segment_share_in_percent": "NA"}


def evaluate_name_and_description(
    segment_input: str,
    segment_info: dict,
    prompt_template2: str,
    client=None,
    model_name=MODEL_NAME, max_tokens=MAX_TOKENS, temperature=TEMPERATURE, top_p=TOP_P, top_k=TOP_K,
) -> str:
    assert client is not None
    prompt2 = prompt_template2.format(input_str=segment_input, output_str=json.dumps(segment_info))
    try:
        response = query_llm(prompt2, model_name=model_name, client=client, max_tokens=max_tokens,
                             temperature=temperature, top_p=top_p, top_k=top_k)
        if response.lower().startswith("y"):
            return "yes"
        return "no"
    except Exception as e:
        print(e)
        return "no"


def search_google_images_selenium(search_query, num_images=10, save_dir=None):
    """
    Search Google Images using Selenium and optionally save the results.
    
    Args:
        search_query (str): The search term or description
        num_images (int): Number of images to retrieve (default: 5)
        save_dir (str): Directory to save images (default: None)
        
    Returns:
        list: List of images in base64
    """
    # Initialize the webdriver (Chrome)
    options = webdriver.ChromeOptions()
    options.add_argument('--headless')  # Run in headless mode
    options.add_argument('--no-sandbox')
    driver = webdriver.Chrome(options=options)

    # Navigate to Google Images
    search_url = f"https://www.google.com/search?q={search_query}&tbm=isch"
    driver.get(search_url)

    # Wait for images to load
    time.sleep(2)

    # Find image elements
    image_elements = driver.find_elements(By.CSS_SELECTOR, "g-img > img")
    base64_binaries = []
    count = 0
    for img_tag in image_elements:
        if count >= num_images:
            break
        try:
            src = img_tag.get_attribute("src")
            if not src or not src.startswith("data:image/"):
                continue
            base64_binary = src.split("base64,")[-1]
            mime_type = src.split(";")[0].split(":")[1]
            file_type = mime_type.split("/")[-1]
            if len(base64_binary) < 5000:
                continue
            if file_type == "gif":
                continue
            base64_binaries.append((file_type, base64_binary))
            
                
            # Save image if directory is specified
            if save_dir:
                if not os.path.exists(save_dir):
                    os.makedirs(save_dir)
                alt_text = img_tag.get_attribute("alt") or "image"
                filename = f"{alt_text}-{count}.{file_type}"
            
                image_binary = base64.b64decode(base64_binary)
                output_path = f"{save_dir}/{filename}"
                with open(output_path, "wb") as f:
                    f.write(image_binary)
            count += 1
        except Exception as e:
            print(f"Error processing image: {str(e)}")
            continue

    driver.quit()
    return base64_binaries


def find_the_best_image(images, search_query, prompt_template=None,
                        max_retries=MAX_RETRIES, 
                        client=None, model_name=MODEL_NAME,
                        max_tokens=MAX_TOKENS, temperature=TEMPERATURE, top_p=TOP_P, top_k=TOP_K,
                        default_image_base64=DEFAULT_IMAGE_BASE64, 
                        default_image_filetype=DEFAULT_IMAGE_FILETYPE):
    assert client is not None
    if prompt_template is None:
        prompt = f"""
You are an assistant to an analyst who is creating a presentation material.
Here we are trying to find the best image for the given description.

<description>
{search_query}
</description>

<guidelines>
1. The image has to have people in it, but do not pick the one with more than 5 people.
2. It has to be a photographed image.
3. Do not pick an inappropriate image because the image will be used in a presentation.
4. Your response should be just one integer. For example, if the 3rd image out of 10 images is picked, the response should be just "3".
</guidelines>
"""
    else:
        prompt = prompt_template.format(search_query=search_query)
    additional_contents = []
    for image in images:
        additional_contents.append(
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": f"image/{image[0]}",
                    "data": image[1]
                }
            }
        )
    num_retries = 0
    while num_retries < max_retries:
        try:
            message = query_llm(prompt, model_name=model_name, client=client, max_tokens=max_tokens, temperature=temperature, top_p=top_p, top_k=top_k,
                               additional_contents=additional_contents)
            best_image_index = int(message) - 1
            return images[best_image_index]
        except Exception as e:
            print(e)
            num_retries += 1
    return default_image_filetype, default_image_base64


def save_base64_image_to_file(
    image_base64: str, 
    filename: str = "image",
    new_size: int = 0,
    new_filetype: str = "png",
) -> str:
    """A function to save a base64-encoded image to a file.

    Args:
        image_base64 (str): base64-encoded image
        filename (str): name of the file without an extension
        new_size (int): new size of the image
        new_filetype (str): new file type

    Returns:
        image file name (str)
    """
    image_binary = base64.b64decode(image_base64)
    img = PIL.Image.open(io.BytesIO(image_binary))
    if new_size:  # crop and resize to a squared image if new_size is given
        image_width, image_height = img._size
        if image_width != image_height:
            if image_width > image_height:
                left_crop = (image_width - image_height) // 2
                right_crop = left_crop + image_height
                top_crop = 0
                bottom_crop = image_height
                image_size = image_height
            else:
                left_crop = (image_width - image_height) // 2
                right_crop = left_crop + image_height
                top_crop = 0
                bottom_crop = image_height
                image_size = image_width
            img = img.crop((left_crop, top_crop, right_crop, bottom_crop))
        else:
            image_size = image_width
        if image_size != new_size:
            img = img.resize((new_size, new_size), PIL.Image.LANCZOS)
    new_filename = f"{filename}.{new_filetype}"
    img.save(new_filename, format=new_filetype)
    return new_filename


def generate_slide(
    segments_df: pd.DataFrame,
    num_segments: int = 4,
    template_pptx_prefix: str = ".",
    image_type: str = "png",
    final_pptx: str = "segments.pptx"
):
    if template_pptx_prefix.startswith("gs://") or final_pptx.startswith("gs://"):
        storage_client = storage.Client()
    # (1) Read template pptx
    template_pptx = f"{template_pptx_prefix}/template{num_segments}.pptx"
    filename_split = template_pptx.rsplit("/", maxsplit=1)
    if len(filename_split) == 1:
        path = "."
        filename = filename_split[0]
    elif len(filename_split) == 2:
        path, filename = filename_split
    else:
        raise Exception(f"Check file, {template_pptx}.")
    if path.startswith("gs://"):
        filename_split2 = path[5:].split("/", maxsplit=1)
        bucket = storage_client.bucket(filename_split2[0])
        if len(filename_split2) == 1:
            blob = bucket.blob(filename)
        else:
            blob = bucket.blob(filename_split2[1] + "/" + filename)
        blob.download_to_filename(filename)
        prs = Presentation(filename)
    else:
        prs = Presentation(template_pptx)

    # (2) Update the slide
    segments_df["segment_id"] = segments_df["segment_id"].astype(int)
    slide = prs.slides[0]
    shapes = slide.shapes
    table = shapes[-1 - num_segments].table
    for segment_id in range(1, num_segments+1):
        segment_name, segment_description, segment_share = segments_df.loc[segments_df["segment_id"] == segment_id].iloc[0][
            ["segment_name", "segment_description", "segment_share_in_percent"]]
        # segment name
        row = 0
        tf = table.cell(row, segment_id-1).text_frame
        p = tf.paragraphs[0]
        p.text = segment_name
        p.font.bold = True
        p.font.size = Pt(14)
        # segment image
        image_filename = f"best_image{segment_id}.{image_type}"
        imgPic = shapes[-num_segments+segment_id-1]._pic
        imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
        imgPart = slide.part.related_part(imgRID)
        with open(image_filename, 'rb') as f:
            rImgBlob = f.read()
            imgPart._blob = rImgBlob
        # segment share
        row = 2
        tf = table.cell(row, segment_id-1).text_frame
        p = tf.paragraphs[0]
        p.text = f"{segment_share}% of customers" 
        p.font.size = Pt(12)
    
        # segment description
        row = 3
        tf = table.cell(row, segment_id-1).text_frame
        p = tf.paragraphs[1]
        p.text = segment_description
        p.font.italic = True
        p.font.size = Pt(12)
    # (3) Write the slide 
    filename_split = final_pptx.rsplit("/", maxsplit=1)
    if len(filename_split) == 1:
        path = "."
        filename = filename_split[0]
    elif len(filename_split) == 2:
        path, filename = filename_split
    else:
        raise Exception(f"Check file, {final_pptx}.")
    if path.startswith("gs://"):
        prs.save(filename)
        filename_split2 = path[5:].split("/", maxsplit=1)
        bucket = storage_client.bucket(filename_split2[0])
        if len(filename_split2) == 1:
            blob = bucket.blob(filename)
        else:
            blob = bucket.blob(filename_split2[1] + "/" + filename)
        blob.upload_from_filename(filename)
    else:
        prs.save(path + "/" + filename)
    return final_pptx